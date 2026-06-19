"""
Text and metadata extraction from files.
Returns a structured result used by Tiers 3-5 for embedding and LLM context.
All path operations use the \\?\ long-path prefix.
"""
import os
import struct
from dataclasses import dataclass, field
from pathlib import Path

LONG_PATH_PREFIX = "\\\\?\\"


def long_path(p: str | Path) -> str:
    s = str(p)
    if s.startswith(LONG_PATH_PREFIX):
        return s
    if s.startswith("\\\\"):
        return "\\\\?\\UNC\\" + s[2:]
    return LONG_PATH_PREFIX + s


@dataclass
class ExtractionResult:
    text: str = ""
    metadata: dict = field(default_factory=dict)
    exif: dict | None = None
    thumbnail: bytes | None = None
    error: str | None = None

    @property
    def snippet(self) -> str:
        return self.text[:500] if self.text else ""


def extract(file_path: str | Path) -> ExtractionResult:
    path = Path(file_path)
    lp = long_path(path)
    ext = path.suffix.lower().lstrip(".")

    try:
        if not os.path.exists(lp):
            return ExtractionResult(error="File not found")
        if os.path.getsize(lp) == 0:
            return ExtractionResult(error="Zero-byte file")
    except OSError as e:
        return ExtractionResult(error=str(e))

    extractors = {
        "pdf":  _extract_pdf,
        "docx": _extract_docx,
        "doc":  _extract_docx,
        "xlsx": _extract_xlsx,
        "xls":  _extract_xlsx,
        "pptx": _extract_pptx,
        "ppt":  _extract_pptx,
        "lnk":  _extract_lnk,
    }

    image_exts = {"jpg", "jpeg", "heic", "heif", "png", "gif", "bmp", "tiff", "tif", "webp", "cr2", "nef", "arw"}
    text_exts = {"txt", "md", "rst", "csv", "json", "yaml", "yml", "toml", "xml", "html", "htm", "css", "log", "ini", "cfg"}

    if ext in extractors:
        return extractors[ext](lp, path)
    if ext in image_exts:
        return _extract_image(lp, path)
    if ext in text_exts:
        return _extract_text(lp, path)

    return ExtractionResult()


def _extract_pdf(lp: str, path: Path) -> ExtractionResult:
    try:
        import pypdf
        reader = pypdf.PdfReader(lp)
        pages = []
        for page in reader.pages[:10]:
            try:
                pages.append(page.extract_text() or "")
            except Exception:
                pass
        text = "\n".join(pages)
        meta = {}
        if reader.metadata:
            meta = {k.lstrip("/"): str(v) for k, v in reader.metadata.items() if v}
        return ExtractionResult(text=text.strip(), metadata=meta)
    except Exception as e:
        return ExtractionResult(error=f"pdf: {e}")


def _extract_docx(lp: str, path: Path) -> ExtractionResult:
    try:
        import docx
        doc = docx.Document(lp)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        meta = {}
        props = doc.core_properties
        for attr in ("author", "title", "subject", "keywords", "created", "modified"):
            val = getattr(props, attr, None)
            if val:
                meta[attr] = str(val)
        return ExtractionResult(text=text, metadata=meta)
    except Exception as e:
        return ExtractionResult(error=f"docx: {e}")


def _extract_xlsx(lp: str, path: Path) -> ExtractionResult:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(lp, read_only=True, data_only=True)
        lines = []
        for sheet in wb.worksheets[:3]:
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append("\t".join(cells))
                row_count += 1
                if row_count >= 50:
                    break
        wb.close()
        return ExtractionResult(text="\n".join(lines))
    except Exception as e:
        return ExtractionResult(error=f"xlsx: {e}")


def _extract_pptx(lp: str, path: Path) -> ExtractionResult:
    try:
        from pptx import Presentation
        prs = Presentation(lp)
        lines = []
        for slide in prs.slides[:20]:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
        return ExtractionResult(text="\n".join(lines))
    except Exception as e:
        return ExtractionResult(error=f"pptx: {e}")


def _extract_lnk(lp: str, path: Path) -> ExtractionResult:
    try:
        import LnkParse3
        with open(lp, "rb") as f:
            lnk = LnkParse3.lnk_file(f)
        info = lnk.get_json()
        target = info.get("link_info", {}).get("local_base_path", "")
        return ExtractionResult(
            text=target,
            metadata={"target": target, "lnk_name": path.name},
        )
    except Exception as e:
        return ExtractionResult(error=f"lnk: {e}")


def _extract_image(lp: str, path: Path) -> ExtractionResult:
    try:
        from PIL import Image
        from PIL.ExifTags import TAGS, GPSTAGS
        img = Image.open(lp)
        exif_data: dict = {}
        thumbnail: bytes | None = None

        raw_exif = img._getexif() if hasattr(img, "_getexif") else None
        if raw_exif:
            for tag_id, value in raw_exif.items():
                tag = TAGS.get(tag_id, str(tag_id))
                if isinstance(value, bytes):
                    continue
                exif_data[tag] = str(value)

            gps_info = raw_exif.get(34853)  # GPSInfo tag
            if gps_info:
                gps = {}
                for k, v in gps_info.items():
                    gps[GPSTAGS.get(k, k)] = v
                location = _reverse_geocode(gps)
                if location:
                    exif_data["_location"] = location

        try:
            img.thumbnail((256, 256))
            import io
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=70)
            thumbnail = buf.getvalue()
        except Exception:
            pass

        img.close()

        text_parts = []
        for key in ("Make", "Model", "DateTimeOriginal", "DateTime", "_location"):
            if key in exif_data:
                text_parts.append(f"{key}: {exif_data[key]}")

        return ExtractionResult(
            text="\n".join(text_parts),
            metadata={"width": img.size[0] if hasattr(img, "size") else None,
                      "height": img.size[1] if hasattr(img, "size") else None},
            exif=exif_data if exif_data else None,
            thumbnail=thumbnail,
        )
    except Exception as e:
        return ExtractionResult(error=f"image: {e}")


def _extract_text(lp: str, path: Path) -> ExtractionResult:
    try:
        import chardet
        with open(lp, "rb") as f:
            raw = f.read(65536)
        detected = chardet.detect(raw)
        encoding = detected.get("encoding") or "utf-8"
        try:
            text = raw.decode(encoding, errors="replace")
        except (LookupError, UnicodeDecodeError):
            text = raw.decode("utf-8", errors="replace")
        return ExtractionResult(text=text)
    except Exception as e:
        return ExtractionResult(error=f"text: {e}")


def _reverse_geocode(gps: dict) -> str | None:
    try:
        lat = _gps_to_decimal(gps.get("GPSLatitude"), gps.get("GPSLatitudeRef", "N"))
        lon = _gps_to_decimal(gps.get("GPSLongitude"), gps.get("GPSLongitudeRef", "E"))
        if lat is None or lon is None:
            return None
        import reverse_geocoder
        results = reverse_geocoder.search([(lat, lon)], verbose=False)
        if results:
            r = results[0]
            return f"{r.get('name', '')}, {r.get('admin1', '')}, {r.get('cc', '')}"
    except Exception:
        pass
    return None


def _gps_to_decimal(coord, ref: str) -> float | None:
    if coord is None:
        return None
    try:
        def to_float(v):
            if isinstance(v, tuple) and len(v) == 2:
                return v[0] / v[1] if v[1] != 0 else 0.0
            return float(v)
        degrees = to_float(coord[0])
        minutes = to_float(coord[1])
        seconds = to_float(coord[2])
        decimal = degrees + minutes / 60 + seconds / 3600
        if ref in ("S", "W"):
            decimal = -decimal
        return decimal
    except Exception:
        return None
